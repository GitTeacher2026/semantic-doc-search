"""Extract text from supported file types, including Office and Arabic filenames."""

from __future__ import annotations

import re
import shutil
import subprocess
import tempfile
from pathlib import Path

SUPPORTED_EXTENSIONS = frozenset(
    {
        ".pdf",
        ".txt",
        ".md",
        ".text",
        ".log",
        ".csv",
        ".doc",
        ".docx",
        ".xls",
        ".xlsx",
        ".ppt",
        ".pptx",
    }
)

EXTENSION_GROUPS: dict[str, set[str]] = {
    "pdf": {".pdf"},
    "word": {".doc", ".docx"},
    "excel": {".xls", ".xlsx"},
    "powerpoint": {".ppt", ".pptx"},
    "text": {".txt", ".md", ".text", ".log", ".csv"},
}

GROUP_LABELS_AR = {
    "pdf": "PDF",
    "word": "Word",
    "excel": "Excel",
    "powerpoint": "PowerPoint",
    "text": "نص",
    "other": "أخرى",
}

GROUP_ICONS = {
    "pdf": "📄",
    "word": "📝",
    "excel": "📊",
    "powerpoint": "📽️",
    "text": "📃",
    "other": "📁",
}


def file_group(suffix: str) -> str:
    lowered = suffix.lower()
    for group, extensions in EXTENSION_GROUPS.items():
        if lowered in extensions:
            return group
    return "other"


def safe_storage_name(original_name: str, doc_id: str) -> str:
    """Preserve Arabic/Unicode in filenames; strip only unsafe path characters."""
    name = Path(original_name).name
    name = re.sub(r'[\\/:*?"<>|\x00-\x1f]', "_", name).strip()
    if not name or name in {".", ".."}:
        name = "مستند"

    stem = Path(name).stem
    suffix = Path(name).suffix.lower()
    if len(name) > 140:
        name = f"{stem[:120]}{suffix}"

    return f"{doc_id}_{name}"


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    return "\n".join(parts).strip()


def _extract_plain_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore").strip()


def _extract_docx(path: Path) -> str:
    from docx import Document

    document = Document(str(path))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _extract_xlsx(path: Path) -> str:
    import openpyxl

    workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"## {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cells:
                parts.append(" | ".join(cells))
    workbook.close()
    return "\n".join(parts).strip()


def _extract_xls(path: Path) -> str:
    import xlrd

    workbook = xlrd.open_workbook(str(path))
    parts: list[str] = []
    for sheet in workbook.sheets():
        parts.append(f"## {sheet.name}")
        for row_idx in range(sheet.nrows):
            cells = [
                str(sheet.cell_value(row_idx, col_idx)).strip()
                for col_idx in range(sheet.ncols)
                if str(sheet.cell_value(row_idx, col_idx)).strip()
            ]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = [f"## شريحة {index}"]
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            if text.strip():
                slide_parts.append(text.strip())
        if len(slide_parts) > 1:
            parts.extend(slide_parts)
    return "\n".join(parts).strip()


def _extract_with_libreoffice(path: Path) -> str:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise ValueError(
            "ملفات .doc و .ppt القديمة تتطلب LibreOffice. "
            "حوّل الملف إلى docx/pptx أو ثبّت LibreOffice."
        )

    with tempfile.TemporaryDirectory() as tmp_dir:
        result = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "txt:Text",
                "--outdir",
                tmp_dir,
                str(path),
            ],
            capture_output=True,
            text=True,
            check=False,
        )
        if result.returncode != 0:
            raise ValueError("تعذّر تحويل الملف باستخدام LibreOffice.")

        converted = Path(tmp_dir) / f"{path.stem}.txt"
        if not converted.exists():
            matches = list(Path(tmp_dir).glob("*.txt"))
            if not matches:
                raise ValueError("لم يُنتج LibreOffice ملف نص.")
            converted = matches[0]
        return converted.read_text(encoding="utf-8", errors="ignore").strip()


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix in {".txt", ".md", ".text", ".log", ".csv"}:
        return _extract_plain_text(path)
    if suffix == ".docx":
        return _extract_docx(path)
    if suffix == ".doc":
        return _extract_with_libreoffice(path)
    if suffix == ".xlsx":
        return _extract_xlsx(path)
    if suffix == ".xls":
        return _extract_xls(path)
    if suffix == ".pptx":
        return _extract_pptx(path)
    if suffix == ".ppt":
        return _extract_with_libreoffice(path)
    raise ValueError(f"نوع الملف غير مدعوم: {suffix}")
